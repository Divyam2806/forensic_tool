import os
from pathlib import Path

def read_text_file(path: str) -> str:
    """Reads text files (including log, csv, txt) with UTF-8 or Latin-1 fallback."""
    try:
        try:
            with open(path, "r", encoding="utf-8") as f:
                return f.read()
        except UnicodeDecodeError:
            with open(path, "r", encoding="latin-1") as f:
                return f.read()
    except Exception as e:
        return f"Error reading text file: {e}"

def read_pdf_file(path: str) -> str:
    """Extracts text from a PDF file using pypdf."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        text_content = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_content.append(text)
        return "\n".join(text_content)
    except Exception as e:
        return f"Error reading PDF file: {e}"

def read_docx_file(path: str) -> str:
    """Extracts text from a Word document using python-docx."""
    try:
        from docx import Document
        doc = Document(path)
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception as e:
        return f"Error reading DOCX file: {e}"

def read_xlsx_file(path: str) -> str:
    """Extracts text from an Excel spreadsheet cell-by-cell using openpyxl."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        text_content = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            for row in sheet.iter_rows(values_only=True):
                row_str = " ".join([str(cell) for cell in row if cell is not None])
                if row_str.strip():
                    text_content.append(row_str)
        return "\n".join(text_content)
    except Exception as e:
        return f"Error reading XLSX file: {e}"

def read_pptx_file(path: str) -> str:
    """Extracts text from shapes in a PowerPoint presentation using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text)
        return "\n".join(text_content)
    except Exception as e:
        return f"Error reading PPTX file: {e}"

def read_file_content(path: str) -> str:
    """Routes the file to the appropriate reader based on extension."""
    path_obj = Path(path).resolve()
    if not path_obj.exists() or not path_obj.is_file():
        return ""
    
    ext = path_obj.suffix.lower()
    
    if ext in [".txt", ".csv", ".log"]:
        return read_text_file(str(path_obj))
    elif ext == ".pdf":
        return read_pdf_file(str(path_obj))
    elif ext == ".docx":
        return read_docx_file(str(path_obj))
    elif ext == ".xlsx":
        return read_xlsx_file(str(path_obj))
    elif ext == ".pptx":
        return read_pptx_file(str(path_obj))
    else:
        return ""
